"""
UniversalParser — Agent 万能文件解析器

parse_file(path, prompt) 工具，替代 spawn_subagent 做文件分析。
支持图片/文档/音频/视频/网页，根据 prompt 自动路由到最优模型。

架构：
```
parse_file(path, prompt)
  → 格式检测
  → Handler 路由（Image / Doc / Audio / Video / URL）
  → Handler 内小模型决策（要不要切块？要不要 Omni？）
  → 执行 → 返回
```
"""

import os
import re
import json
import base64
import hashlib
import logging
import asyncio
import aiohttp
import shutil
import tempfile
import subprocess
import httpx
from collections import OrderedDict
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple

from services.tools.base import AgentToolsServiceBase
from services.tool_registry import tool

logger = logging.getLogger(__name__)

# ── Qwen 文本对话（OpenAI-compatible 端点） ─────────────────
# dashscope.Generation.call() 对 qwen3.6-flash 路由到错误端点（"url error"），
# 所以走 httpx 直接打 compatible-mode 端点。

QWEN_CHAT_URL = "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions"



async def _qwen_chat(
    model: str,
    messages: list,
    api_key: str,
    temperature: float = 0.3,
    max_tokens: int = 2000,
) -> Optional[str]:
    """Call DashScope via OpenAI-compatible endpoint (true async HTTP via aiohttp).

    Auto-retries on failure with progressively simpler requests.
    Returns the content string, or None on error.
    """
    # Retry strategy: on first failure, retry with shorter max_tokens
    retry_configs = [
        {"max_tokens": max_tokens, "text_ratio": 1.0},   # original
        {"max_tokens": max(max_tokens // 2, 500), "text_ratio": 0.7},  # half tokens, 70% text
        {"max_tokens": max(max_tokens // 4, 200), "text_ratio": 0.4},  # quarter tokens, 40% text
    ]

    for attempt, cfg in enumerate(retry_configs):
        try:
            # Truncate user message text if this is a retry
            msgs = messages
            if attempt > 0:
                msgs = []
                for m in messages:
                    if m["role"] == "user" and len(m["content"]) > 500:
                        trunc_len = int(len(m["content"]) * cfg["text_ratio"])
                        msgs.append({"role": "user", "content": m["content"][:trunc_len]})
                    else:
                        msgs.append(m)

            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=300)) as session:
                async with session.post(
                    QWEN_CHAT_URL,
                    headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                    json={"model": model, "messages": msgs, "temperature": temperature, "max_tokens": cfg["max_tokens"]},
                ) as resp:
                    if resp.status != 200:
                        if attempt < len(retry_configs) - 1:
                            continue
                        return None
                    data = await resp.json()
            choices = data.get("choices", [])
            if choices:
                return choices[0].get("message", {}).get("content")
            return None
        except (asyncio.TimeoutError, aiohttp.ServerTimeoutError, Exception):
            if attempt < len(retry_configs) - 1:
                continue
            return None
    return None


# ── VLM 多模态（图片/视频帧）通过 chat completions + base64 走 qwen3.6-flash ──
# 统一走 httpx chat completions（与 _qwen_chat 同一端点），避免
# MultiModalConversation SDK 的额外依赖和路径问题。

VLM_MODEL = "qwen3.6-flash"


async def _vlm_chat(
    image_paths: list,
    prompt: str,
    api_key: str,
    max_tokens: int = 4096,
) -> Optional[str]:
    """Call qwen3.6-flash with images via httpx (base64 inline).

    image_paths: list of paths to PNG/JPEG files.
    Returns the content string, or None on error.
    """
    try:
        content = []
        for path in image_paths:
            with open(path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("utf-8")
            ext = path.rsplit(".", 1)[-1].lower()
            mime = f"image/{'png' if ext == 'png' else 'jpeg'}"
            content.append({
                "type": "image_url",
                "image_url": {"url": f"data:{mime};base64,{b64}"}
            })
        content.append({"type": "text", "text": prompt})

        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=300)) as session:
            async with session.post(
                QWEN_CHAT_URL,
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={
                    "model": VLM_MODEL,
                    "messages": [{"role": "user", "content": content}],
                    "temperature": 0.3,
                    "max_tokens": max_tokens,
                },
            ) as resp:
                data = await resp.json()
        choices = data.get("choices", [])
        if choices:
            return choices[0].get("message", {}).get("content")
        return None
    except Exception:
        return None

# ── 文件类型检测 ──────────────────────────────────────────────

# 扩展名 → 类别
EXT_CATEGORY = {
    # 图片
    "jpg": "image", "jpeg": "image", "png": "image", "gif": "image",
    "webp": "image", "bmp": "image", "svg": "image", "heic": "image",
    # 纯文本
    "txt": "doc", "md": "doc", "json": "doc", "csv": "doc",
    "yaml": "doc", "yml": "doc", "xml": "doc", "html": "doc", "htm": "doc",
    "log": "doc", "py": "doc", "js": "doc", "ts": "doc", "cpp": "doc",
    "c": "doc", "h": "doc", "java": "doc", "go": "doc", "rs": "doc",
    # 文档
    "pdf": "doc", "docx": "doc", "doc": "doc",
    "pptx": "doc", "ppt": "doc",
    "xlsx": "doc", "xls": "doc",
    # 音频
    "mp3": "audio", "wav": "audio", "m4a": "audio", "ogg": "audio",
    "flac": "audio", "aac": "audio", "wma": "audio",
    # 视频
    "mp4": "video", "mov": "video", "avi": "video", "mkv": "video",
    "webm": "video", "flv": "video",
}

# 可直接当文本读取的扩展名
TEXT_EXTS = {
    "txt", "md", "json", "csv", "yaml", "yml", "xml", "html", "htm",
    "log", "py", "js", "ts", "cpp", "c", "h", "java", "go", "rs",
}

# 大文件阈值（字节）—— 超过则走小模型决策
DOC_LARGE_THRESHOLD = 20 * 1024  # 20KB
# 文本截断长度（避免 LLM context 爆炸）
TEXT_TRUNCATE_CHARS = 100_000

# 平行处理限制
MAX_CONCURRENT = 20     # 最多20路并发
MAX_UNITS = 100         # 最多100个单元

# 支持的扩展名（用于错误提示）
SUPPORTED_EXTS = sorted(EXT_CATEGORY.keys())

# 反向索引：category → ext list（避免每次错误消息 O(n) 扫描）
_EXTS_BY_CAT: Dict[str, List[str]] = {}
for _ext, _cat in EXT_CATEGORY.items():
    _EXTS_BY_CAT.setdefault(_cat, []).append(_ext)
for _lst in _EXTS_BY_CAT.values():
    _lst.sort()


# ── 小模型路由决策（SmartRouter 风格） ─────────────────────────

_SR_CACHE_MAX = 256
_SR_CACHE: "OrderedDict[str, Dict[str, Any]]" = OrderedDict()


def _sr_cache_key(content_type: str, prompt: str, file_size: int, duration_s: Optional[float]) -> str:
    """稳定 cache key（避免 Python hash() 跨进程不一致）"""
    h = hashlib.sha256()
    h.update(content_type.encode("utf-8"))
    h.update(b"\0")
    h.update(prompt.encode("utf-8", errors="replace"))
    h.update(b"\0")
    h.update(str(file_size).encode("utf-8"))
    h.update(b"\0")
    h.update(str(int(duration_s or 0)).encode("utf-8"))
    return h.hexdigest()


def _sr_cache_get(key: str) -> Optional[Dict[str, Any]]:
    if key in _SR_CACHE:
        # LRU：移到末尾
        _SR_CACHE.move_to_end(key)
        return _SR_CACHE[key]
    return None


def _sr_cache_put(key: str, value: Dict[str, Any]) -> None:
    _SR_CACHE[key] = value
    _SR_CACHE.move_to_end(key)
    if len(_SR_CACHE) > _SR_CACHE_MAX:
        _SR_CACHE.popitem(last=False)


async def _route_decision(
    content_type: str,
    prompt: str,
    file_size: int,
    duration_s: Optional[float] = None,
) -> Dict[str, Any]:
    """
    小模型（Qwen-turbo）预判路由策略。

    返回:
    - needs_chunking: bool — 文档是否需要切块检索
    - requires_omni: bool — 音视频是否需要 Omni
    - reasoning: str — 决策理由
    """
    cache_key = _sr_cache_key(content_type, prompt, file_size, duration_s)
    cached = _sr_cache_get(cache_key)
    if cached is not None:
        return cached

    size_desc = f"{file_size / 1024:.0f} KB"
    if duration_s:
        size_desc += f"（{duration_s/60:.1f} 分钟）"

    system_prompt = f"""你是一个文件解析路由决策器。根据用户的问题判断最佳处理方案。

当前文件类型: {content_type}
文件大小: {size_desc}
用户问题: {prompt}

请只返回 JSON，不要多余内容：

1. 对于文档（txt/pdf/docx 等）：
   needs_chunking: 用户问题是否只需要文件中某一部分的信息？
   - 如果问题是宽泛的（"总结"、"讲什么"）→ false（全文即可）
   - 如果是具体的（"第二章"、"某人的观点"、"数据"）→ true（需切块检索）
   - 如果文件很小（<10KB）→ false（全文喂不贵）

2. 对于音频/视频：
   requires_omni: 用户问题是否需要理解声音/画面的「原始感知」？
   - 需要：背景音识别（"海浪声"）、音色情绪（"生气吗"）、画面内容（"在做什么"）
   - 不需要：只问内容/摘要 → false

4. 新模式判断（mode）：
   mode="solve": 用户要求解/做/计算文档中的题目（如"解所有题"、"求答案"）
   mode="grade": 用户要求批改/评分（如"批改答案"、"判断对错"）
   mode="extract_terms": 用户要求提取/列出离散元素（术语、公式等）
   mode="summarize": 用户要求总结/概括文档内容
   mode="default": 其他情况

   has_units: 文档是否包含可拆分处理的离散单元（题、术语等）
   unit_type: "problem"（试题）| "answer_pair"（题目+学生答案）| "term"（术语/公式）
   unit_count_estimate: 预估单元数量

返回示例:
{{"mode":"solve","has_units":true,"unit_type":"problem","unit_count_estimate":12,"needs_chunking":false,"reasoning":"数学试卷，12道选择题"}}

返回 JSON: {{"mode": str, "has_units": bool, "unit_type": str|null, "unit_count_estimate": int, "needs_chunking": bool, "requires_omni": bool, "reasoning": "简短理由"}}
"""

    result: Dict[str, Any]
    try:
        from config import settings

        text = await _qwen_chat(
            model="qwen3.6-flash",
            api_key=settings.QWEN_API_KEY,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt or "总结这个文件"},
            ],
            temperature=0.1,
            max_tokens=300,
        )
        if text:
            # 提取 JSON（处理 ```json 包裹的情况）
            text = re.sub(r"```json\s*|\s*```", "", text).strip()
            # 容错：找最外层 { ... }
            m = re.search(r"\{[\s\S]*\}", text)
            if m:
                parsed = json.loads(m.group(0))
                result = {
                    "mode": str(parsed.get("mode", "default")),
                    "has_units": bool(parsed.get("has_units", False)),
                    "unit_type": parsed.get("unit_type") if parsed.get("unit_type") in ("problem", "answer_pair", "term") else None,
                    "unit_count_estimate": int(parsed.get("unit_count_estimate", 0) or 0),
                    "needs_chunking": bool(parsed.get("needs_chunking", False)),
                    "requires_omni": bool(parsed.get("requires_omni", False)),
                    "reasoning": str(parsed.get("reasoning", "")),
                }
            else:
                result = _sr_fallback(content_type, "no JSON in output")
        else:
            result = _sr_fallback(content_type, "无输出")
    except Exception as e:
        logger.warning(f"SR decision failed: {e}")
        result = _sr_fallback(content_type, str(e))

    _sr_cache_put(cache_key, result)
    return result


def _sr_fallback(content_type: str, reason: str) -> Dict[str, Any]:
    """SR 失败时的保守默认值"""
    if content_type in ("audio", "video"):
        return {
            "mode": "default",
            "has_units": False,
            "unit_type": None,
            "unit_count_estimate": 0,
            "needs_chunking": False,
            "requires_omni": True,
            "reasoning": f"fallback: {reason}",
        }
    return {
        "mode": "default",
        "has_units": False,
        "unit_type": None,
        "unit_count_estimate": 0,
        "needs_chunking": False,
        "requires_omni": False,
        "reasoning": f"fallback: {reason}",
    }


# ── 平行单元处理（提取 → 并发 → 聚合） ─────────────────────────


async def _extract_units(
    text: str,
    mode: str,
    unit_type: Optional[str],
    prompt_hint: str,
) -> List[Dict[str, Any]]:
    """用一次 LLM 调用从文档中提取可并行处理的单元列表。

    Returns: list of {"id": int, "stem": str, "metadata": dict}
    """
    SYSTEM_PROMPTS = {
        "solve": """你是一个文档解析器。从文档中逐一提取所有独立的题目。
返回 JSON 数组，每项包含：
- id: 题号（从1开始）
- stem: 完整的题目内容（包括题干、选项、配图描述等所有信息）
- metadata: { "type": "choice"|"fill"|"解答"|"其他" }

注意：
- 每题独立，不要遗漏
- 选择题要保留所有选项
- 解答题保留全部文字
- 即使题目之间有共享材料（如同一段阅读材料），也单独列出""",
        "grade": """你是一个批改助手。从文档中提取所有包含"题目"和对应"学生答案"的对。

返回 JSON 数组，每项包含：
- id: 题号
- stem: 题目内容
- student_answer: 学生给出的答案
- metadata: {}""",
        "extract_terms": """你是一个术语提取器。从文档中提取所有离散的信息单元。

返回 JSON 数组，每项包含：
- id: 序号
- stem: 术语/公式/标题 + 内容
- metadata: { "type": "term"|"formula"|"table"|"section" }""",
    }

    try:
        from config import settings
    except Exception:
        logger.warning("config import failed in _extract_units")
        return []

    prompt = SYSTEM_PROMPTS.get(mode, SYSTEM_PROMPTS["solve"])

    # Retry with progressively shorter text if first attempt fails
    text_lengths = [80000, 50000, 20000, 8000]
    for attempt, max_chars in enumerate(text_lengths):
        if attempt > 0:
            logger.info(f"_extract_units retry {attempt+1}, truncating to {max_chars} chars")

        text_resp = await _qwen_chat(
            model="qwen3.6-flash",
            api_key=settings.QWEN_API_KEY,
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": f"文档内容：\n\n{text[:max_chars]}\n\n用户问题：{prompt_hint}"},
            ],
            temperature=0.1,
            max_tokens=4096 if attempt == 0 else 2048,
        )
        if text_resp:
            break
    else:
        return []

    text_resp = re.sub(r"```json\s*|\s*```", "", text_resp).strip()
    # 非贪婪匹配：取第一个完整的 JSON 数组（模型有时在数组后追加解释文字）
    m = re.search(r"\[[\s\S]*?\]", text_resp)
    if not m:
        return []

    try:
        parsed = json.loads(m.group(0))
    except Exception as e:
        logger.warning(f"_extract_units JSON parse failed: {e}")
        return []

    units: List[Dict[str, Any]] = []
    for i, item in enumerate(parsed):
        if not isinstance(item, dict):
            continue
        stem = item.get("stem", "")
        if not stem:
            continue
        units.append({
            "id": item.get("id", i + 1),
            "stem": stem,
            "metadata": item.get("metadata", {}) if isinstance(item.get("metadata"), dict) else {},
        })
        if mode == "grade" and item.get("student_answer"):
            units[-1]["student_answer"] = item["student_answer"]

    return units


async def _parallel_process(
    units: List[Dict[str, Any]],
    mode: str,
    prompt_hint: str,
    page_images: Optional[List[str]] = None,
) -> List[Dict[str, Any]]:
    """并发处理所有单元，返回结果列表。"""

    if len(units) > MAX_UNITS:
        units = units[:MAX_UNITS]

    try:
        from config import settings
    except Exception:
        logger.warning("config import failed in _parallel_process")
        return [{"id": u["id"], "result": "（配置加载失败）"} for u in units]

    def _build_task_prompt(unit: Dict[str, Any], mode: str) -> str:
        """为单个单元构建处理 prompt"""
        if mode == "solve":
            return f"""你是一个解题助手。请解答以下题目，给出详细的解题步骤和最终答案。

题目：
{unit['stem']}

请按以下格式输出：
## 解题过程
（写出详细步骤）
## 最终答案
（写出最终答案）"""

        elif mode == "grade":
            return f"""你是一个批改助手。判断答案是否正确，如有错误请指出并给出正确答案。

题目：
{unit['stem']}

学生答案：
{unit.get('student_answer', '（无学生答案）')}

请按以下格式输出：
## 判断
正确/错误（需改正）/部分正确
## 评语
...
## 正确答案
..."""


        elif mode == "extract_terms":
            return f"""请解释以下内容：

{unit['stem']}

请给出简要的解释说明。"""

        else:
            return f"""请处理以下内容：

{unit['stem']}

给出你的分析和结论。"""

    sem = asyncio.Semaphore(MAX_CONCURRENT)

    async def _process_one(unit: Dict[str, Any]) -> Dict[str, Any]:
        """处理单个单元（带超时保护）"""
        async with sem:
            task_prompt = _build_task_prompt(unit, mode)
            try:
                # If we have page images and are in solve mode, use multimodal
                if page_images and mode == "solve":
                    content = []
                    for img_path in page_images:
                        with open(img_path, "rb") as f:
                            b64 = base64.b64encode(f.read()).decode("utf-8")
                        content.append({
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{b64}"}
                        })
                    content.append({"type": "text", "text": task_prompt})
                    messages = [
                        {"role": "system", "content": "你是一个专业的解题和分析助手。"},
                        {"role": "user", "content": content},
                    ]
                    text = await asyncio.wait_for(
                        _qwen_chat(
                            model="qwen3.6-flash",
                            api_key=settings.QWEN_API_KEY,
                            messages=messages,
                            temperature=0.3,
                            max_tokens=2048,
                        ),
                        timeout=240,
                    )
                else:
                    text = await asyncio.wait_for(
                        _qwen_chat(
                            model="qwen3.6-flash",
                            api_key=settings.QWEN_API_KEY,
                            messages=[
                                {"role": "system", "content": "你是一个专业的解题和分析助手。"},
                                {"role": "user", "content": task_prompt},
                            ],
                            temperature=0.3,
                            max_tokens=2048,
                        ),
                        timeout=240,  # 每单元最大240s
                    )
                return {"id": unit["id"], "result": text or "（无结果）"}
            except asyncio.TimeoutError:
                return {"id": unit["id"], "result": "（处理超时，已跳过）"}

    results = await asyncio.gather(
        *[_process_one(u) for u in units],
        return_exceptions=True,
    )

    processed: List[Dict[str, Any]] = []
    for i, r in enumerate(results):
        if isinstance(r, Exception):
            processed.append({"id": units[i]["id"], "result": f"（处理失败：{r}）"})
        else:
            processed.append(r)

    return processed


async def _aggregate(results: List[Dict[str, Any]], mode: str) -> str:
    """聚合所有并行处理的结果"""
    if not results:
        return "（无处理结果）"

    if mode == "solve":
        lines: List[str] = []
        for r in results:
            lines.append(f"## 第 {r['id']} 题\n{r['result']}\n")
        return "\n".join(lines)

    # For other modes, use LLM to summarize
    summary_input = json.dumps(
        [{"id": r["id"], "result": r["result"][:500]} for r in results],
        ensure_ascii=False, indent=2,
    )
    try:
        from config import settings
        text = await _qwen_chat(
            model="qwen3.6-flash",
            api_key=settings.QWEN_API_KEY,
            messages=[
                {"role": "system", "content": "你是一个结果汇总助手。汇总以下并行处理的结果，生成清晰的结构化报告。"},
                {"role": "user", "content": f"请汇总以下并行处理结果：\n\n{summary_input}"},
            ],
            max_tokens=3000,
        )
        return text or "（汇总失败）"
    except Exception as e:
        logger.warning(f"_aggregate failed: {e}")
        return "（汇总失败）"


# ── LLM 文本回答 ──────────────────────────────────────────────

async def _llm_answer(
    prompt: str,
    context: str,
    max_tokens: int = 2000,
    page_images: Optional[List[str]] = None,
) -> str:
    """用 Qwen-turbo 基于上下文回答"""
    try:
        from config import settings

        # If we have page images, use VLM instead
        if page_images:
            text = await _vlm_chat(
                page_images,
                f"文件内容：\n\n{context[:TEXT_TRUNCATE_CHARS]}\n\n用户问题：{prompt}",
                settings.QWEN_API_KEY,
                max_tokens=max_tokens,
            )
            if text:
                return text
            return "（模型未返回结果）"

        text = await _qwen_chat(
            model="qwen3.6-flash",
            api_key=settings.QWEN_API_KEY,
            messages=[
                {"role": "system", "content": "你是一个专业的文件分析助手。基于提供的文件内容回答用户的问题。"},
                {"role": "user", "content": f"文件内容：\n\n{context[:TEXT_TRUNCATE_CHARS]}\n\n用户问题：{prompt}"},
            ],
            temperature=0.3,
            max_tokens=max_tokens,
        )
        if text:
            return text
        return "（模型未返回结果）"
    except Exception as e:
        return f"（回答时出错：{e}）"


# ── Omni 调用（音频理解；ASR 也走 Omni） ─────────────────────

# BLOCKER: 真正的 ASR 走 DashScope Paraformer (services/tools/...
# 中可加 _asr_transcribe)，但当前 API key 权限只允许 Omni 多模态，
# 所以音频转写也通过 qwen3.5-omni-plus 完成。授权问题解决后可在
# 这里加分支：短音频用 Paraformer（便宜），长音频用 Omni。


async def _omni_analyze(
    audio_path: str,
    prompt: str,
    video_url: Optional[str] = None,
    extra_text_hint: str = "",
) -> str:
    """
    调用 Qwen3.5-Omni-Plus 分析音频或视频。
    音频输入支持本地文件路径（SDK 自动上传）。
    """
    try:
        from dashscope import MultiModalConversation
        from config import settings

        # 反幻觉 prompt：明确告诉模型"只听音频"避免它凭训练数据瞎编
        anti_hallucination = (
            "这是一段新的音频内容，请只基于你从音频波形中直接听到的信息回答，"
            "不要依赖你的训练数据知识。如果你没有听清，请明确说'无法听清'。"
        )
        text_parts = [anti_hallucination]
        if extra_text_hint:
            text_parts.append(extra_text_hint)
        text_parts.append(prompt or "请描述这段音频的内容")

        messages = [{"role": "user", "content": []}]
        if video_url:
            messages[0]["content"].append(
                {"type": "video_url", "video_url": {"url": video_url}}
            )
        else:
            messages[0]["content"].append({"audio": audio_path})
        messages[0]["content"].append({"text": "\n".join(text_parts)})

        response = MultiModalConversation.call(
            model="qwen3.5-omni-plus",
            api_key=settings.QWEN_API_KEY,
            messages=messages,
            result_format="message",
        )

        output = response.get("output", {})
        choices = output.get("choices", [])
        text = ""
        if choices:
            content = choices[0].get("message", {}).get("content", [])
            for c in content:
                if isinstance(c, dict) and "text" in c:
                    text += c["text"]
        return text or "（Omni 未返回内容）"
    except Exception as e:
        return f"（Omni 分析时出错：{e}）"


# ── 文本提取（文档类，二进制格式本地解析） ─────────────────────

def _extract_text(path: str, max_chars: int = TEXT_TRUNCATE_CHARS) -> str:
    """从本地文件提取文本（PDF/DOCX/PPTX/XLSX/TXT/...）"""
    ext = Path(path).suffix.lower().lstrip(".")

    try:
        if ext in TEXT_EXTS:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read(max_chars)

        if ext == "pdf":
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(path)
                parts: List[str] = []
                total = 0
                for page in doc:
                    page_text = page.get_text()
                    parts.append(page_text)
                    total += len(page_text)
                    if total > max_chars:
                        break
                return "".join(parts)[:max_chars]
            except ImportError:
                return "（PyMuPDF 未安装）"
            except Exception as e:
                return f"（PDF 提取失败：{e}）"

        if ext == "docx":
            try:
                import docx
                doc = docx.Document(path)
                return "\n".join(p.text for p in doc.paragraphs)[:max_chars]
            except ImportError:
                return "（python-docx 未安装）"
            except Exception as e:
                return f"（DOCX 提取失败：{e}）"

        if ext in ("pptx", "ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(path)
                parts = []
                total = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                            total += len(shape.text)
                            if total > max_chars:
                                break
                    if total > max_chars:
                        break
                return "\n".join(parts)[:max_chars]
            except ImportError:
                return "（python-pptx 未安装）"
            except Exception as e:
                return f"（PPTX 提取失败：{e}）"

        if ext in ("xlsx", "xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                parts = []
                total = 0
                for sheet in wb.worksheets:
                    parts.append(f"=== Sheet: {sheet.title} ===")
                    total += len(parts[-1])
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join(str(c) if c is not None else "" for c in row)
                        if row_text.strip():
                            parts.append(row_text)
                            total += len(row_text)
                            if total > max_chars:
                                break
                    if total > max_chars:
                        break
                return "\n".join(parts)[:max_chars]
            except ImportError:
                return "（openpyxl 未安装）"
            except Exception as e:
                return f"（XLSX 提取失败：{e}）"

        return f"（不支持的文档格式: {ext}）"

    except Exception as e:
        return f"（文本提取时出错：{e}）"


# ── ffmpeg / ffprobe 工具 ─────────────────────────────────────

def _convert_to_wav(path: str) -> Optional[str]:
    """将音频/视频音频轨转 16kHz mono WAV；失败返回 None"""
    wav_path = path + ".wav"
    try:
        result = subprocess.run(
            ["ffmpeg", "-y", "-i", path,
             "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1",
             wav_path],
            capture_output=True, text=True, timeout=300
        )
        if result.returncode == 0 and os.path.exists(wav_path) and os.path.getsize(wav_path) > 0:
            return wav_path
        # 清理可能产生的零字节文件
        if os.path.exists(wav_path):
            try:
                os.remove(wav_path)
            except OSError:
                pass
        return None
    except FileNotFoundError:
        logger.warning("ffmpeg 未安装，跳过 WAV 转换")
        return None
    except subprocess.TimeoutExpired:
        logger.warning(f"ffmpeg convert 超时: {path}")
        return None
    except Exception as e:
        logger.warning(f"ffmpeg convert failed: {e}")
        return None


def _get_duration(path: str) -> float:
    """获取音频/视频时长（秒）"""
    try:
        result = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", path],
            capture_output=True, text=True, timeout=30
        )
        s = result.stdout.strip()
        return float(s) if s else 0.0
    except Exception:
        return 0.0


def _extract_keyframes(video_path: str, interval: int = 10, max_frames: int = 30) -> Tuple[List[str], str]:
    """从视频中按 interval 秒抽关键帧（jpg 列表），返回 (frames, frame_dir)"""
    duration = _get_duration(video_path)
    if duration <= 0:
        return [], ""
    frame_dir = tempfile.mkdtemp(prefix="vf_frames_")
    try:
        # 一次 ffmpeg 调用：fps 过滤器按 interval 秒抽帧，远快于 30 次串行调用
        cmd = [
            "ffmpeg", "-y", "-i", video_path,
            "-vf", f"fps=1/{interval}",
            "-frames:v", str(max_frames),
            "-q:v", "2",
            os.path.join(frame_dir, "frame_%04d.jpg"),
        ]
        subprocess.run(cmd, capture_output=True, timeout=300)
        frames = sorted(
            os.path.join(frame_dir, f)
            for f in os.listdir(frame_dir)
            if f.startswith("frame_") and f.endswith(".jpg")
        )
        # 过滤掉零字节文件
        frames = [f for f in frames if os.path.getsize(f) > 0]
        return frames, frame_dir
    except Exception as e:
        logger.warning(f"Keyframe extraction failed: {e}")
        return [], frame_dir


def _frames_to_vlm_content(frames: List[str], header_text: str = "请按时间顺序描述这些视频画面中发生了什么。") -> List[dict]:
    """将帧图片打包成 VLM 可用的 content 列表"""
    content: List[dict] = [{"text": header_text}]
    for fp in frames:
        try:
            with open(fp, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("ascii")
            content.append({"image": f"data:image/jpeg;base64,{b64}"})
        except Exception as e:
            logger.warning(f"frame {fp} read failed: {e}")
    return content


def _cleanup_paths(*paths: Optional[str]) -> None:
    """清理一组可能的临时文件路径（不存在的忽略）"""
    for p in paths:
        if p and os.path.exists(p):
            try:
                os.remove(p)
            except OSError:
                pass


# ── VLM 图片分析 ─────────────────────────────────────────────

def _extract_mm_text(response: dict, fallback: str = "（模型未返回内容）") -> str:
    """从 DashScope MultiModalConversation 返回值中提取纯文本"""
    choices = response.get("output", {}).get("choices", [])
    if not choices:
        return fallback
    text = ""
    for c in choices[0].get("message", {}).get("content", []):
        if isinstance(c, dict) and "text" in c:
            text += c["text"]
    return text or fallback


async def _vlm_analyze_image(image_path: str, prompt: str) -> str:
    """Qwen VLM 分析图片（qwen3.6-flash via httpx）"""
    try:
        from config import settings

        text = await _vlm_chat([image_path], prompt, settings.QWEN_API_KEY)
        return text or "（VLM 未返回内容）"
    except Exception as e:
        return f"（VLM 分析时出错：{e}）"


async def _vlm_analyze_content(content: List[dict]) -> str:
    """Qwen VLM 分析多模态 content（用于视频关键帧，qwen3.6-flash via httpx）。

    content: list of {"image": "data:image/...;base64,..."} or {"text": str}
    """
    try:
        from config import settings
        import base64 as _b64

        image_paths: List[str] = []
        text_parts: List[str] = []
        for c in content:
            if "image" in c:
                # data URI → 写临时文件给 _vlm_chat
                img = c["image"]
                if isinstance(img, str) and img.startswith("data:image"):
                    # 解析 "data:image/jpeg;base64,XXX"
                    try:
                        header, b64data = img.split(",", 1)
                        mime = header.split(";")[0].split(":", 1)[1]  # "image/jpeg"
                        ext = "png" if "png" in mime else "jpg"
                        tmp = os.path.join(
                            tempfile.gettempdir(),
                            f"parse_vlm_{os.urandom(4).hex()}.{ext}",
                        )
                        with open(tmp, "wb") as f:
                            f.write(_b64.b64decode(b64data))
                        image_paths.append(tmp)
                    except Exception:
                        continue
                elif isinstance(img, str):
                    image_paths.append(img)
            if "text" in c:
                text_parts.append(c["text"])

        prompt = "\n".join(text_parts) if text_parts else "请描述这些图片的内容。"
        text = await _vlm_chat(image_paths, prompt, settings.QWEN_API_KEY)
        return text or "（VLM 未返回内容）"
    except Exception as e:
        return f"（VLM 帧分析失败：{e}）"


# ── 核心 Mixin ─────────────────────────────────────────────

class ParseFileMixin(AgentToolsServiceBase):
    """parse_file 工具——万能文件解析器"""

    @tool(
        description=(
            "万能文件解析器：分析任何文件（图片/文档/音频/视频/网页），"
            "回答你对文件内容的任何问题。会根据文件类型自动选择最合适的模型。\n\n"
            "支持的文件类型：\n"
            "- 图片: jpg, jpeg, png, gif, webp, svg, bmp, heic（→ Qwen VL-Max）\n"
            "- 文档: pdf, docx, pptx, xlsx, txt, md, json, csv, yaml, xml, html, py 等（→ 文本提取 + Qwen-Turbo）\n"
            "- 音频: mp3, wav, m4a, ogg, flac, aac, wma（→ Qwen3.5-Omni-Plus）\n"
            "- 视频: mp4, mov, avi, mkv, webm, flv（→ 关键帧 VLM + 音频 Omni）\n"
            "- 网址: http://, https://（→ 抓取 + Qwen-Turbo）\n\n"
            "提示：对于大文件，具体的问题（\u201c第二章讲什么\u201d）比笼统的问题（\u201c总结一下\u201d）结果更精确。"
        ),
        category="file",
    )
    async def parse_file(self, path: str, prompt: str) -> str:
        """
        万能文件解析器。

        :param path: 文件在 VFS 中的路径（如 /workspace/photo.jpg），或 http(s) URL
        :param prompt: 你想问关于这个文件的问题（如 "这是什么？"、"总结一下"）
        """
        try:
            return await self._do_parse(path, prompt)
        except Exception as e:
            logger.exception("parse_file error: %s", e)
            return f"（ERROR: 解析时出错：{e}）"

    # ── 主调度 ─────────────────────────────────────────────

    async def _do_parse(self, path: str, prompt: str) -> str:
        ext = Path(path).suffix.lower().lstrip(".")
        category = EXT_CATEGORY.get(ext, "unknown")

        if category == "unknown" and re.match(r"^https?://", path):
            category = "url"

        if category == "unknown":
            return self._unsupported_msg(ext)

        if category == "image":
            return await self._handle_image(path, prompt)
        if category == "doc":
            return await self._handle_doc(path, prompt)
        if category == "audio":
            return await self._handle_audio(path, prompt)
        if category == "video":
            return await self._handle_video(path, prompt)
        if category == "url":
            return await self._handle_url(path, prompt)
        # 不可达：上面 if/elif 链已穷尽 EXT_CATEGORY 的所有值
        return self._unsupported_msg(ext)  # pragma: no cover

    def _unsupported_msg(self, ext: str) -> str:
        return (
            f"（不支持的文件格式: .{ext or '(空)'}。"
            f"当前支持：图片({', '.join(_EXTS_BY_CAT.get('image', []))})、"
            f"文档({', '.join(_EXTS_BY_CAT.get('doc', []))})、"
            f"音频({', '.join(_EXTS_BY_CAT.get('audio', []))})、"
            f"视频({', '.join(_EXTS_BY_CAT.get('video', []))})、"
            f"网址(http/https)）"
        )

    # ── Image Handler ──────────────────────────────────────

    async def _handle_image(self, path: str, prompt: str) -> str:
        """图片 → Qwen VL-Max"""
        file_bytes = await self._vfs_read_bytes(path)
        if not file_bytes:
            return "（找不到文件或无法读取）"

        suffix = Path(path).suffix or ".img"
        tmp_path = os.path.join(tempfile.gettempdir(), f"parse_img_{os.urandom(4).hex()}{suffix}")
        try:
            with open(tmp_path, "wb") as f:
                f.write(file_bytes)
            return await _vlm_analyze_image(tmp_path, prompt)
        except Exception as e:
            return f"（处理图片失败：{e}）"
        finally:
            _cleanup_paths(tmp_path)

    # ── Doc Handler ─────────────────────────────────────────

    async def _handle_doc(self, path: str, prompt: str) -> str:
        """文档 → 文本提取 → 小模型决策（要不要切块） → LLM 回答"""
        ext = Path(path).suffix.lower().lstrip(".")
        is_text = ext in TEXT_EXTS
        _page_images: List[str] = []  # hoisted for use after else block

        # 文本文件走 VFS；二进制格式（pdf/docx/pptx/xlsx）下载到本地
        if is_text:
            try:
                text = await self._vfs.async_read_file(path)
            except Exception as e:
                return f"（读取文件失败：{e}）"
            if text.startswith("Error:"):
                return text
        else:
            file_bytes = await self._vfs_read_bytes(path)
            if not file_bytes:
                return "（找不到文件或无法读取）"
            tmp_path = os.path.join(tempfile.gettempdir(), f"parse_doc_{os.urandom(4).hex()}.{ext}")
            try:
                with open(tmp_path, "wb") as f:
                    f.write(file_bytes)
                text = _extract_text(tmp_path)

                # VLM fallback for PDFs: garbled text > 10% OR embedded meaningful images
                needs_vlm = False
                if ext == "pdf":
                    doc = fitz.open(tmp_path)

                    # Smarter image detection: skip full-page backgrounds (>90%) and tiny icons (<5%)
                    has_meaningful = False
                    page_rects = [page.rect for page in doc]
                    for i, page in enumerate(doc):
                        pw, ph = page_rects[i].width, page_rects[i].height
                        for img in page.get_images(full=True):
                            bbox = page.get_image_bbox(img)
                            if bbox is None or bbox.is_empty:
                                continue
                            rw = bbox.width / pw
                            rh = bbox.height / ph
                            # Skip full-page background (>90%) and tiny icons (<5%)
                            if rw > 0.9 and rh > 0.9:
                                continue
                            if rw < 0.05 and rh < 0.05:
                                continue
                            has_meaningful = True
                            break
                        if has_meaningful:
                            break
                    doc.close()

                    if has_meaningful:
                        needs_vlm = True
                        logger.info("PDF has meaningful embedded images, using VLM OCR")
                    elif text and text.count("\ufffd") / max(len(text), 1) > 0.10:
                        needs_vlm = True
                        garbled_count = text.count("\ufffd")
                        logger.info(
                            f"PDF garbled {garbled_count}/{len(text)} chars "
                            f"({100*garbled_count/max(len(text),1):.1f}%), using VLM OCR"
                        )
                        text = ""

                    if needs_vlm:
                        doc = fitz.open(tmp_path)
                        for i in range(len(doc)):
                            pix = doc[i].get_pixmap(dpi=200)
                            img_path = os.path.join(
                                tempfile.gettempdir(),
                                f"parse_pdf_page_{os.urandom(4).hex()}_{i}.png",
                            )
                            pix.save(img_path)
                            _page_images.append(img_path)
                        doc.close()
                    if _page_images:
                        from config import settings
                        prompt_text = f"""请识别这份文档的全部内容，包括所有文字和数学公式（用LaTeX格式输出）。注意保留：
- 所有题目编号和题干
- 数学公式（用 $...$ 或 $$...$$ 包围）
- 几何图形中的标注文字
- 任何图表或表格中的数据
- 用户的问题：{prompt}"""
                        text = await _vlm_chat(_page_images, prompt_text, settings.QWEN_API_KEY) or ""
            except Exception as e:
                return f"（读取文件失败：{e}）"
            finally:
                _cleanup_paths(tmp_path)

        if not text or not text.strip():
            _cleanup_paths(*_page_images)
            return "（文件内容为空或无法提取文字）"

        file_size = len(text.encode("utf-8"))

        # 小文件直接全文
        if file_size <= DOC_LARGE_THRESHOLD:
            result = await _llm_answer(prompt, text[:TEXT_TRUNCATE_CHARS], page_images=_page_images)
            _cleanup_paths(*_page_images)
            return result

        # 大文件：小模型决策要不要切块
        decision = await _route_decision("doc", prompt, file_size)
        logger.info(f"Doc SR decision: {decision}")

        mode = decision.get("mode", "default")

        # 平行单元处理模式
        if mode in ("solve", "grade", "extract_terms") and decision.get("has_units"):
            logger.info(
                f"Parallel mode: {mode}, estimated units: {decision.get('unit_count_estimate', '?')}"
            )
            try:
                units = await _extract_units(
                    text, mode, decision.get("unit_type"), prompt
                )
                if units:
                    logger.info(f"Extracted {len(units)} units for parallel processing")
                    results = await _parallel_process(units, mode, prompt, page_images=_page_images)
                    _cleanup_paths(*_page_images)
                    return await _aggregate(results, mode)
                logger.warning("extract_units returned empty, falling back")
            except Exception as e:
                logger.warning(f"Parallel processing failed: {e}, falling back")

        # 智能切块检索：SemanticChunker 切分 → embedding 相似度 → top-K
        if decision.get("needs_chunking") and file_size > 50000:
            try:
                from services.semantic_chunker import chunk as sem_chunk
                from services.embedding_service import EmbeddingService

                chunks = await sem_chunk(text)
                l1_chunks = [c for c in chunks if c.level == 1]

                if l1_chunks:
                    embedder = EmbeddingService()
                    prompt_emb = await embedder.embed(prompt)
                    chunk_embs = await embedder.embed_batch([c.text for c in l1_chunks])

                    import math
                    def _cosim(a, b):
                        if not a or not b or len(a) != len(b):
                            return 0.0
                        dot = sum(x * y for x, y in zip(a, b))
                        na = math.sqrt(sum(x * x for x in a))
                        nb = math.sqrt(sum(x * x for x in b))
                        return dot / (na * nb) if na and nb else 0.0

                    scored = [(_cosim(prompt_emb, ce), c) for ce, c in zip(chunk_embs, l1_chunks)]
                    scored.sort(key=lambda x: x[0], reverse=True)

                    TOP_K = 5
                    MIN_SCORE = 0.3
                    top_chunks = [c for score, c in scored[:TOP_K] if score > MIN_SCORE]

                    if top_chunks:
                        context = "\n\n".join(c.text for c in top_chunks)
                        # 附加 L2 超块：包含 top chunk 起始 50 字符的更大上下文（最多 2 个）
                        l2_chunks = [c for c in chunks if c.level == 2]
                        relevant_l2 = [
                            c for c in l2_chunks
                            if any(c.text.find(tc.text[:50]) >= 0 for tc in top_chunks)
                        ]
                        if relevant_l2:
                            context += "\n\n【扩展上下文】\n\n" + "\n\n".join(
                                c.text for c in relevant_l2[:2]
                            )
                    else:
                        # 没有任何 chunk 越过阈值 → 取前 5 个 L1 做兜底
                        context = "\n\n".join(c.text for c in l1_chunks[:5])

                    context = context[:TEXT_TRUNCATE_CHARS]  # 安全截断
                    result = await _llm_answer(prompt, context, page_images=_page_images)
                    _cleanup_paths(*_page_images)
                    return result
            except Exception as e:
                logger.warning(f"SemanticChunker 路径失败，回退到截断：{e}")

        # 默认：直接截断
        truncated = text[:TEXT_TRUNCATE_CHARS]
        result = await _llm_answer(prompt, truncated, page_images=_page_images)
        _cleanup_paths(*_page_images)
        return result

    # ── Audio Handler ──────────────────────────────────────

    async def _handle_audio(self, path: str, prompt: str) -> str:
        """音频 → 16kHz mono WAV → Qwen3.5-Omni-Plus
        (ASR 也走 Omni，授权问题解决后可在此加 Paraformer 分支)"""
        file_bytes = await self._vfs_read_bytes(path)
        if not file_bytes:
            return "（找不到文件或无法读取）"

        suffix = Path(path).suffix or ".audio"
        tmp_path = os.path.join(tempfile.gettempdir(), f"parse_audio_{os.urandom(4).hex()}{suffix}")
        wav_path: Optional[str] = None
        try:
            with open(tmp_path, "wb") as f:
                f.write(file_bytes)

            # 优先转 16kHz mono WAV（Omni 对此最稳定）；失败时回退到原始文件
            wav_path = _convert_to_wav(tmp_path) or tmp_path

            extra_hint = (
                "如果用户要求转写/听写，请把音频中实际说出的文字逐字记录下来；"
                "无法听清的部分明确标注'（听不清）'。"
            )
            return await _omni_analyze(
                audio_path=wav_path,
                prompt=prompt,
                extra_text_hint=extra_hint,
            )
        except Exception as e:
            return f"（处理音频失败：{e}）"
        finally:
            _cleanup_paths(tmp_path, wav_path if wav_path and wav_path != tmp_path else None)

    # ── Video Handler ──────────────────────────────────────

    async def _handle_video(self, path: str, prompt: str) -> str:
        """视频 → 抽音频（Omni 理解）+ 抽关键帧（VLM 描述）→ LLM 综合回答"""
        file_bytes = await self._vfs_read_bytes(path)
        if not file_bytes:
            return "（找不到文件或无法读取）"

        suffix = Path(path).suffix or ".video"
        tmp_path = os.path.join(tempfile.gettempdir(), f"parse_video_{os.urandom(4).hex()}{suffix}")
        audio_path: Optional[str] = None
        frame_paths: List[str] = []
        frame_dir: str = ""
        try:
            with open(tmp_path, "wb") as f:
                f.write(file_bytes)

            duration = _get_duration(tmp_path)
            logger.info(f"Video parse start: duration={duration:.1f}s, size={os.path.getsize(tmp_path)}")

            # 抽音频和关键帧在主线程顺序做（ffmpeg 多线程反而慢）
            audio_path = _convert_to_wav(tmp_path)
            frame_paths, frame_dir = _extract_keyframes(tmp_path, interval=10, max_frames=30)

            # 音频 Omni + 关键帧 VLM 并行：两者独立，可省 ~50% 延迟
            async def _analyze_audio() -> str:
                if not audio_path or not os.path.exists(audio_path):
                    return ""
                return await _omni_analyze(
                    audio_path=audio_path,
                    prompt="请仔细听这段音频，逐字转写其中的语音内容，并描述重要的非语音声音（背景音、音效等）。",
                    extra_text_hint="重点：只基于听到的内容回答，不要编造。",
                )

            async def _analyze_frames() -> str:
                if not frame_paths:
                    return ""
                vlm_content = _frames_to_vlm_content(
                    frame_paths,
                    header_text="以下是按时间顺序抽取的视频关键帧，请按顺序描述每帧画面内容及变化。",
                )
                return await _vlm_analyze_content(vlm_content)

            audio_desc, frame_desc = await asyncio.gather(_analyze_audio(), _analyze_frames())

            if not audio_desc and not frame_desc:
                return "（视频解析失败：未能提取任何画面或音频信息）"

            # 给画面和音频各分一半预算，避免其中一项把另一项挤出
            half = TEXT_TRUNCATE_CHARS // 2
            context_parts: List[str] = []
            if frame_desc:
                context_parts.append(f"【画面描述】\n{frame_desc[:half]}")
            if audio_desc:
                context_parts.append(f"【音频内容】\n{audio_desc[:half]}")
            context_parts.append(f"用户问题：{prompt}")
            return await _llm_answer(prompt, "\n\n".join(context_parts))

        except Exception as e:
            return f"（处理视频失败：{e}）"
        finally:
            extras = [audio_path] if audio_path and audio_path != tmp_path else []
            _cleanup_paths(tmp_path, *extras)
            if frame_dir and os.path.isdir(frame_dir):
                shutil.rmtree(frame_dir, ignore_errors=True)

    # ── URL Handler ────────────────────────────────────────

    async def _handle_url(self, path: str, prompt: str) -> str:
        """网页 → httpx 抓取 → 去标签 → LLM 回答"""
        import httpx
        try:
            resp = httpx.get(path, timeout=15, follow_redirects=True)
            resp.raise_for_status()
        except Exception as e:
            return f"（抓取网页失败：{e}）"

        # HTML 去标签
        html = resp.text[:TEXT_TRUNCATE_CHARS * 2]  # 多取一些，标签清理后会缩水
        html = re.sub(r"<script[^>]*>[\s\S]*?</script>", " ", html, flags=re.IGNORECASE)
        html = re.sub(r"<style[^>]*>[\s\S]*?</style>", " ", html, flags=re.IGNORECASE)
        html = re.sub(r"<[^>]+>", " ", html)
        text = re.sub(r"\s+", " ", html).strip()[:TEXT_TRUNCATE_CHARS]

        if not text:
            return "（页面内容为空或无法抓取）"

        return await _llm_answer(prompt, text)

    # ── VFS 工具方法 ────────────────────────────────────────

    async def _vfs_read_bytes(self, path: str) -> Optional[bytes]:
        """从 VFS 读取二进制文件（走 StorageService 异步接口）"""
        try:
            resolved = self._resolve(path)
        except Exception as e:
            logger.warning(f"_resolve failed for {path}: {e}")
            return None
        if not resolved:
            return None
        try:
            return await self.storage.get_file_content_async(resolved)
        except Exception as e:
            logger.warning(f"_vfs_read_bytes failed for {path}: {e}")
            return None
