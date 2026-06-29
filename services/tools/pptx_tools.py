"""
PPTX Agent 工具 — HTML 幻灯片 → PowerPoint

Agent 把 HTML 幻灯片写入 VFS 文件 → create_pptx 读取 VFS 路径 → GX-Alex (Chromium) 转换 → 保存到 VFS → 返回分享链接.

需要服务器安装：
- node.js
- playwright chromium (for GX-Alex html2pptx)
"""

import os
import sys
import uuid
import logging
import tempfile
import subprocess
from pathlib import Path

from services.tools.base import AgentToolsServiceBase, tool

logger = logging.getLogger(__name__)

# html2pptx 脚本路径
_SCRIPT_DIR = Path(__file__).resolve().parent.parent.parent / "scripts"
_HTML2PPTX_SCRIPT = _SCRIPT_DIR / "html2pptx" / "html2pptx.py"
_CHROMIUM_PATH = os.environ.get(
    "CHROME",
    os.path.expanduser("~/.cache/ms-playwright/chromium-1223/chrome-linux64/chrome"),
)
if not os.path.exists(str(_CHROMIUM_PATH)):
    _playwright_dir = Path.home() / ".cache" / "ms-playwright"
    if _playwright_dir.exists():
        _chromium_dirs = sorted(_playwright_dir.glob("chromium-*/chrome-linux64/chrome"))
        if _chromium_dirs:
            _CHROMIUM_PATH = str(_chromium_dirs[-1])


class PptxToolsMixin(AgentToolsServiceBase):
    """PPTX 生成工具 Mixin"""

    @tool(description=(
        "将 VFS 中的 HTML 幻灯片转换为可编辑的 PowerPoint 文件（.pptx）。"
        "用法：先把 HTML 写到一个 VFS 文件（如 workspace/slides/deck.html），"
        "然后调用此工具传入 vfs_path。"
        "HTML 格式详见 /public/feclaw/tools/create_pptx.md。\n\n"
        "要点：每张幻灯片必须是 <section class=\"deck-slide\">，1280x720px，样式写在 <style> 中。"
        "背景用纯色（如 background: #ffffff），不要用 CSS 渐变（gradient）。\n\n"
        "raster_mode 参数：截图为图片嵌入 PPT，而非可编辑 SVG。"
        "注意 raster_mode=True 时：① 无法做 PPT 动画，② 文字不可编辑。如非必要不要开启。"
    ), category="file")
    def create_pptx(
        self,
        vfs_path: str,
        output_filename: str = "presentation.pptx",
        raster_mode: bool = False,
    ) -> str:
        """
        将 VFS 中的 HTML 幻灯片转换为 PPTX

        Args:
            vfs_path: VFS 文件路径，如 "/workspace/slides/deck.html"
            output_filename: 输出文件名（默认 presentation.pptx）
            raster_mode: 截图模式。True=走 Chromium 截图嵌入 PPT（保留渐变等复杂样式），
                        False=走 GX-Alex SVG 原生转换（PPT 可编辑但渐变简化）。
                        截图模式缺点：① 无法做 PPT 动画，② 文字不可编辑。

        Returns:
            成功返回 VFS 路径和分享链接，失败返回错误信息
        """
        try:
            # 1. 从 VFS 读取 HTML
            html_content = self.vfs.cat(vfs_path)
            if html_content.startswith("Error:"):
                return f"Error: 读取文件失败 ({vfs_path}): {html_content}"

            # 2. 写入临时 HTML 文件
            html_path = os.path.join(
                tempfile.gettempdir(),
                f"pptx_input_{uuid.uuid4().hex[:8]}.html",
            )
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(html_content)

            # 3. 确定输出路径
            output_path = os.path.join(
                tempfile.gettempdir(),
                f"pptx_output_{uuid.uuid4().hex[:8]}.pptx",
            )

            # 4. 根据模式选择转换方式
            if raster_mode:
                success = self._render_raster(html_path, output_path)
            else:
                success = self._render_native(html_path, output_path)

            if not success:
                return "Error: PPTX 转换失败"

            if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                return "Error: PPTX 生成失败（输出为空）"

            # 5. 读取 PPTX 内容并保存到 VFS
            with open(output_path, "rb") as f:
                pptx_bytes = f.read()

            vfs_output_path = f"workspace/presentations/{output_filename}"
            abs_key = f"feclaw/agents/{self.agent_hash}/{vfs_output_path}"
            from services.storage_service import StorageService
            StorageService().upload_file(pptx_bytes, abs_key)

            # 6. 生成分享链接
            share_result = None
            try:
                from services.share_service import create_share_link
                share_result = create_share_link(
                    vfs_path=f"/{vfs_output_path}",
                    mode="share",
                    user_id=self.user_id,
                    expires_hours=0,
                    agent_hash=self.agent_hash,
                )
            except Exception as e:
                logger.warning(f"Share link generation failed: {e}")

            # 7. 清理临时文件
            try:
                os.unlink(html_path)
                os.unlink(output_path)
            except OSError:
                pass

            # 8. 返回结果
            mode_label = "截图模式" if raster_mode else "原生 SVG 模式"
            result_parts = [
                f"✅ PPTX 已生成（{mode_label}）并保存到 {vfs_output_path}",
                f"大小: {len(pptx_bytes) / 1024:.0f} KB",
            ]
            if share_result:
                result_parts.append(f"分享链接: {share_result['url']}")
            return "\n".join(result_parts)

        except subprocess.TimeoutExpired:
            return "Error: PPTX 转换超时（>120 秒）"
        except Exception as e:
            logger.warning(f"create_pptx failed: {e}")
            return f"Error: PPTX 生成失败: {str(e)[:200]}"

    def _render_native(self, html_path: str, output_path: str) -> bool:
        """原生 SVG 模式：GX-Alex html2pptx + Chromium → 可编辑 PPTX"""
        script = _HTML2PPTX_SCRIPT
        if not script.exists():
            logger.error("html2pptx script not found")
            return False

        chrome = str(_CHROMIUM_PATH)
        if not os.path.exists(chrome):
            logger.error("Chromium not found")
            return False

        cmd = [
            sys.executable or "python3", str(script),
            str(html_path), "-o", str(output_path),
            "--chrome", chrome,
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            logger.warning(f"native render failed: {result.stderr[:500]}")
            return False
        return True

    def _render_raster(self, html_path: str, output_path: str) -> bool:
        """截图模式：Playwright Chromium → PNG 截图 → python-pptx 嵌入"""
        import shutil
        script = str(_SCRIPT_DIR / "pptx-rasterizer.js")
        if not os.path.exists(script):
            logger.error("pptx-rasterizer.js not found")
            return False

        chrome = str(_CHROMIUM_PATH)
        slides_dir = os.path.join(tempfile.gettempdir(), f"pptx_slides_{uuid.uuid4().hex[:8]}")
        os.makedirs(slides_dir, exist_ok=True)

        # 截图
        env = os.environ.copy()
        env["CHROME"] = chrome
        result = subprocess.run(
            ["node", script, "--input", str(html_path), "--output", slides_dir],
            capture_output=True, text=True, timeout=120, env=env,
        )
        if result.returncode != 0:
            logger.warning(f"raster screenshot failed: {result.stderr[:500]}")
            shutil.rmtree(slides_dir, ignore_errors=True)
            return False

        # python-pptx 嵌入截图
        from pptx import Presentation
        from pptx.util import Inches, Emu
        prs = Presentation()
        prs.slide_width = Emu(12192000)   # 1280px @ 96dpi
        prs.slide_height = Emu(6858000)   # 720px @ 96dpi

        png_files = sorted(
            f for f in os.listdir(slides_dir) if f.endswith(".png")
        )
        if not png_files:
            shutil.rmtree(slides_dir, ignore_errors=True)
            return False

        for png_name in png_files:
            png_path = os.path.join(slides_dir, png_name)
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            slide.shapes.add_picture(
                png_path, Emu(0), Emu(0),
                width=Emu(12192000), height=Emu(6858000),
            )

        prs.save(output_path)
        shutil.rmtree(slides_dir, ignore_errors=True)
        return True
